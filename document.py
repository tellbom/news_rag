from flask import Flask, request, jsonify
import logging
import os
import io
import tempfile
import uuid
import numpy as np
from PIL import Image
from io import BytesIO
from typing import Dict, List, Tuple, Optional, Union, Any
import mammoth
import PyPDF2
from paddleocr import PaddleOCR
from bs4 import BeautifulSoup
import requests
from pptx import Presentation
import fitz  # PyMuPDF
import subprocess
import fnmatch

app = Flask(__name__)

# 配置日志记录
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler("document_processor.log")
    ]
)
logger = logging.getLogger(__name__)


class DocumentProcessor:
    """文档处理工具类，用于处理不同类型的文档和图像"""

    def __init__(self, ocr_language='ch', use_angle_cls=True):
        """
        初始化文档处理器

        Args:
            ocr_language: OCR识别的语言，默认为中文
            use_angle_cls: 是否使用方向分类器，用于处理旋转的文本
        """
        # 处理numpy兼容性问题
        if not hasattr(np, 'int'):
            np.int = np.int32

        # 初始化OCR模型
        try:
            self.ocr = PaddleOCR(
                det_model_dir="/root/.paddleocr/whl/det/ch/ch_PP-OCRv3_det_infer",
                rec_model_dir="/root/.paddleocr/whl/rec/ch/ch_PP-OCRv3_rec_infer",
                cls_model_dir="/root/.paddleocr/whl/cls/ch_ppocr_mobile_v2.0_cls_infer",
                use_angle_cls=use_angle_cls,
                lang=ocr_language
            )
            logger.info("OCR模型初始化成功")
        except Exception as e:
            self.ocr = None
            logger.error(f"初始化OCR模型失败: {str(e)}")

        # 配置阈值
        self.text_quality_threshold = {
            'min_text_length': 50,  # 最小文本长度
            'min_text_density': 0.01,  # 最小文本密度
            'expected_non_ascii_ratio': 0.2,  # 预期的非ASCII字符比例（用于中文文档）
            'min_word_confidence': 0.6,  # OCR最小词汇置信度
        }

    def render_ppt_as_images(self, ppt_path, output_dir):
        """
        将PPT直接转换为一系列图像，绕过PDF转换

        Args:
            ppt_path: PPT文件路径
            output_dir: 输出图像的目录

        Returns:
            List[str]: 所有图像文件的路径列表
        """
        try:
            # 创建输出目录
            os.makedirs(output_dir, exist_ok=True)

            # 先尝试将PPT转换为PDF，然后再处理PDF
            # 这个方法比直接转换为图像更可靠
            basename = os.path.basename(ppt_path)
            name_without_ext = os.path.splitext(basename)[0]
            pdf_path = os.path.join(output_dir, f"{name_without_ext}.pdf")

            # 第一步：转换为PDF
            cmd_to_pdf = [
                'libreoffice', '--headless', '--norestore', '--convert-to', 'pdf',
                '--outdir', output_dir, ppt_path
            ]

            logger.info(f"执行命令: {' '.join(cmd_to_pdf)}")
            result_pdf = subprocess.run(cmd_to_pdf, capture_output=True, text=True)

            if result_pdf.returncode != 0 or not os.path.exists(pdf_path):
                logger.error(f"PPT转PDF失败: {result_pdf.stderr}")
                # 尝试使用另一种方式直接转换
                alternative_cmd = [
                    'libreoffice', '--headless', '--norestore',
                    '--convert-to', 'png',
                    '--outdir', output_dir, ppt_path
                ]
                logger.info(f"尝试替代命令: {' '.join(alternative_cmd)}")
                alt_result = subprocess.run(alternative_cmd, capture_output=True, text=True)

                # 查找生成的PNG文件
                image_paths = []
                for file in os.listdir(output_dir):
                    if file.startswith(name_without_ext) and file.endswith('.png'):
                        image_paths.append(os.path.join(output_dir, file))

                if not image_paths:
                    logger.error(f"替代方法也失败: {alt_result.stderr}")
                    # 列出输出目录的内容，帮助调试
                    logger.info(f"输出目录内容: {os.listdir(output_dir)}")
                    return []

                # 按文件名排序
                image_paths.sort()
                logger.info(f"使用替代方法找到 {len(image_paths)} 张图像")
                return image_paths

            # 第二步：从PDF生成图像
            pdf_document = fitz.open(pdf_path)
            image_paths = []

            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x提高分辨率
                image_path = os.path.join(output_dir, f"{name_without_ext}-{page_num + 1}.png")
                pix.save(image_path)
                image_paths.append(image_path)

            pdf_document.close()

            if not image_paths:
                logger.warning("未找到生成的幻灯片图像")
            else:
                logger.info(f"成功生成 {len(image_paths)} 张幻灯片图像")

            return image_paths

        except Exception as e:
            logger.error(f"PPT渲染为图像时出错: {str(e)}")
            # 详细记录异常堆栈，帮助调试
            import traceback
            logger.error(traceback.format_exc())
            return []

    def render_slide_to_image(self,ppt_path, slide_number, output_dir):
        """
        使用LibreOffice将PPT的特定幻灯片渲染为图片

        Args:
            ppt_path: PPT文件路径
            slide_number: 要渲染的幻灯片编号（从0开始）
            output_dir: 输出目录

        Returns:
            str: 渲染后的图片路径，如果渲染失败则返回None
        """
        try:
            # 创建输出目录（如果不存在）
            os.makedirs(output_dir, exist_ok=True)
            output_file = os.path.join(output_dir, f"slide_{slide_number + 1}.png")

            # 使用LibreOffice将指定幻灯片导出为PNG
            cmd = [
                'libreoffice', '--headless', '--convert-to', 'png',
                '--outdir', output_dir, ppt_path,
                f'--export-slide-number={slide_number}'
            ]

            result = subprocess.run(cmd, capture_output=True, text=True)

            if result.returncode != 0:
                logger.error(f"渲染幻灯片失败：{result.stderr}")
                return None

            # 检查是否生成了图片
            if os.path.exists(output_file):
                return output_file
            else:
                # 如果导出的文件名与预期不同，查找该目录下最新的PNG文件
                png_files = [f for f in os.listdir(output_dir) if f.endswith('.png')]
                if png_files:
                    latest_file = max(png_files, key=lambda f: os.path.getctime(os.path.join(output_dir, f)))
                    return os.path.join(output_dir, latest_file)
                else:
                    return None
        except Exception as e:
            logger.error(f"渲染幻灯片时出错: {str(e)}")
            return None

    def sanitize_html(self, html_content: str) -> str:
        """
        清理HTML内容，提取纯文本

        Args:
            html_content: HTML字符串

        Returns:
            str: 从HTML中提取的干净文本
        """
        try:
            soup = BeautifulSoup(html_content, 'html.parser')

            # 移除脚本和样式元素
            for script in soup(["script", "style"]):
                script.extract()

            # 获取文本
            text = soup.get_text(separator=' ', strip=True)

            # 清理空白区域
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)

            return text
        except Exception as e:
            logger.error(f"HTML清理错误: {str(e)}")
            # 如果清理失败，返回原始内容
            return html_content

    def process_image(self, image_data: bytes) -> Tuple[str, float]:
        """
        使用OCR从图像中提取文本

        Args:
            image_data: 二进制图像数据

        Returns:
            Tuple[str, float]: 提取的文本和平均置信度分数
        """
        if self.ocr is None:
            logger.error("OCR模型未初始化")
            return "OCR模型不可用", 0.0

        try:
            # 将二进制数据打开为PIL Image，并转换为RGB三通道
            img = Image.open(BytesIO(image_data)).convert("RGB")
            img_np = np.array(img)

            # 调用OCR
            result = self.ocr.ocr(img_np, cls=True)

            # 提取文字与置信度
            text_parts = []
            confidence_sum = 0.0
            count = 0
            for line in result or []:
                if isinstance(line, list) and line:
                    # 取最后一个检测结果
                    content, score = line[-1]
                    confidence = float(score)
                    if content and confidence > self.text_quality_threshold['min_word_confidence']:
                        text_parts.append(content)
                        confidence_sum += confidence
                        count += 1

            # 计算平均置信度
            avg_confidence = confidence_sum / count if count > 0 else 0.0
            full_text = "\n".join(text_parts).strip()

            if not full_text:
                return "图像中未检测到文本", 0.0

            return full_text, avg_confidence

        except Exception as e:
            logger.error(f"图像处理错误: {e}")
            return f"处理图像失败: {e}", 0.0

    def process_word_document(self, docx_data: bytes) -> str:
        """
        从Word文档中提取文本

        Args:
            docx_data: 二进制Word文档数据

        Returns:
            str: 提取的文本
        """
        try:
            result = mammoth.extract_raw_text(BytesIO(docx_data))
            return result.value
        except Exception as e:
            logger.error(f"Word文档处理错误: {str(e)}")
            return f"处理Word文档失败: {str(e)}"

    def need_ocr(self, extracted_text: str, page_size: tuple) -> bool:
        """
        判断提取的文本是否需要OCR处理

        Args:
            extracted_text: 从文档中提取的文本
            page_size: 页面尺寸 (width, height)

        Returns:
            bool: 如果需要OCR则返回True
        """
        # 如果提取的文本为空或几乎为空
        if len(extracted_text) < self.text_quality_threshold['min_text_length']:
            logger.info(
                f"文本长度不足 ({len(extracted_text)} < {self.text_quality_threshold['min_text_length']}), 需要OCR")
            return True

        # 计算文本密度
        text_density = len(extracted_text) / (page_size[0] * page_size[1])
        if text_density < self.text_quality_threshold['min_text_density']:
            logger.info(
                f"文本密度不足 ({text_density:.4f} < {self.text_quality_threshold['min_text_density']}), 需要OCR")
            return True

        # 检查非ASCII字符比例（针对中文文档）
        non_ascii_chars = sum(1 for c in extracted_text if ord(c) > 127)
        non_ascii_ratio = non_ascii_chars / max(len(extracted_text), 1)

        # 如果文档有大量内容，但非ASCII字符（如中文）比例过低，可能是未正确提取
        if non_ascii_ratio < self.text_quality_threshold['expected_non_ascii_ratio']:
            logger.info(
                f"非ASCII字符比例过低 ({non_ascii_ratio:.4f} < {self.text_quality_threshold['expected_non_ascii_ratio']}), 需要OCR")
            return True

        # 简单语义检查 - 检测明显的乱码
        # 例如连续的特殊字符或无意义的重复模式
        weird_chars_count = sum(1 for c in extracted_text if ord(c) < 32 or ord(c) > 126 and ord(c) < 256)
        weird_ratio = weird_chars_count / max(len(extracted_text), 1)
        if weird_ratio > 0.1:  # 超过10%的奇怪字符
            logger.info(f"检测到可能的乱码 (奇怪字符比例: {weird_ratio:.4f}), 需要OCR")
            return True

        return False

    def process_pdf_hybrid(self, pdf_data: bytes) -> Dict[str, Any]:
        """
        使用混合方法处理PDF，根据每页内容质量决定使用文本提取或OCR

        Args:
            pdf_data: 二进制PDF数据

        Returns:
            Dict: 包含提取的文本和处理元数据的字典
        """
        try:
            # 首先尝试打开PDF文件
            pdf_file = BytesIO(pdf_data)

            # 使用PyMuPDF (fitz)处理PDF，它对于文本提取比PyPDF2更好
            pdf_document = fitz.open(stream=pdf_file, filetype="pdf")
            total_pages = len(pdf_document)

            logger.info(f"PDF共有{total_pages}页")

            # 存储处理结果
            result = {
                "text": "",
                "metadata": {
                    "total_pages": total_pages,
                    "pages_extraction_method": {},
                    "processing_summary": {}
                }
            }

            pages_text = []
            pages_using_ocr = 0
            pages_using_extraction = 0

            # 逐页处理
            for page_num in range(total_pages):
                page = pdf_document[page_num]
                page_size = (page.rect.width, page.rect.height)

                # 尝试提取文本
                extracted_text = page.get_text()

                # 判断是否需要OCR
                if self.need_ocr(extracted_text, page_size):
                    # 将PDF页面渲染为图像
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x提高分辨率，提高OCR质量
                    img_data = pix.tobytes("png")

                    # 对图像进行OCR
                    ocr_text, confidence = self.process_image(img_data)
                    pages_text.append(ocr_text)

                    # 记录处理方法
                    result["metadata"]["pages_extraction_method"][page_num + 1] = {
                        "method": "ocr",
                        "confidence": confidence
                    }
                    pages_using_ocr += 1

                    logger.info(f"第{page_num + 1}页使用OCR处理，置信度: {confidence:.2f}")
                else:
                    # 使用文本提取结果
                    pages_text.append(extracted_text)

                    # 记录处理方法
                    result["metadata"]["pages_extraction_method"][page_num + 1] = {
                        "method": "text_extraction"
                    }
                    pages_using_extraction += 1

                    logger.info(f"第{page_num + 1}页使用文本提取")

            # 合并所有页面的文本
            result["text"] = "\n\n".join(pages_text)

            # 添加处理摘要
            result["metadata"]["processing_summary"] = {
                "pages_using_ocr": pages_using_ocr,
                "pages_using_extraction": pages_using_extraction,
                "ocr_percentage": pages_using_ocr / total_pages * 100 if total_pages > 0 else 0
            }

            return result

        except Exception as e:
            logger.error(f"PDF混合处理错误: {str(e)}")
            return {"text": f"处理PDF文档失败: {str(e)}", "metadata": {"error": str(e)}}

    def process_ppt_hybrid(self, ppt_data: bytes, is_pptx: bool = True) -> Dict[str, Any]:
        """
        使用混合方法处理PPT文件，提取文本并对需要的部分使用OCR

        Args:
            ppt_data: 二进制PPT/PPTX数据
            is_pptx: 是否是PPTX格式（True）或PPT格式（False）

        Returns:
            Dict: 包含提取的文本和处理元数据的字典
        """
        try:
            # 创建临时目录用于处理文件
            with tempfile.TemporaryDirectory() as temp_dir:
                # 创建临时PPT文件
                ppt_ext = '.pptx' if is_pptx else '.ppt'
                temp_ppt_path = os.path.join(temp_dir, f"presentation{ppt_ext}")
                with open(temp_ppt_path, 'wb') as temp_file:
                    temp_file.write(ppt_data)

                # 使用python-pptx提取文本内容
                prs = Presentation(temp_ppt_path)
                total_slides = len(prs.slides)

                logger.info(f"PPT共有{total_slides}张幻灯片")

                # 存储处理结果
                result = {
                    "text": "",
                    "metadata": {
                        "total_slides": total_slides,
                        "slides_extraction_method": {},
                        "processing_summary": {}
                    }
                }

                # 初始化统计数据
                slides_text = []
                slides_using_ocr = 0
                slides_using_extraction = 0
                slides_with_images = 0

                # 步骤1: 直接将PPT渲染为图像（不通过PDF）
                image_dir = os.path.join(temp_dir, "images")
                os.makedirs(image_dir, exist_ok=True)
                slide_images = self.render_ppt_as_images(temp_ppt_path, image_dir)
                images_available = len(slide_images) > 0

                if images_available:
                    logger.info(f"成功将PPT渲染为图像，共 {len(slide_images)} 张")
                else:
                    logger.warning("PPT渲染为图像失败，将仅使用文本提取")

                # 步骤2: 逐张幻灯片处理
                for slide_num, slide in enumerate(prs.slides):
                    # 2.1: 提取幻灯片文本
                    slide_text = []

                    # 提取幻灯片标题
                    if slide.shapes.title and slide.shapes.title.text:
                        slide_text.append(f"标题: {slide.shapes.title.text}")

                    # 提取幻灯片中的所有文本
                    text_content = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text_content.append(shape.text)

                    # 合并幻灯片文本内容
                    extracted_text = "\n".join(text_content)

                    # 2.2: 检查是否有图片或文本质量不足
                    has_images = any(shape.shape_type == 13 for shape in slide.shapes)  # 13 = 图片
                    if has_images:
                        slides_with_images += 1

                    # 估算幻灯片尺寸
                    slide_size = (prs.slide_width, prs.slide_height)
                    needs_ocr = self.need_ocr(extracted_text, slide_size) or has_images

                    # 2.3: 处理逻辑
                    # 如果需要OCR且有可用的图像，则使用OCR
                    if needs_ocr and images_available and slide_num < len(slide_images):
                        try:
                            # 读取对应幻灯片的图像
                            with open(slide_images[slide_num], 'rb') as img_file:
                                slide_image_data = img_file.read()

                            # 对图像进行OCR
                            ocr_text, confidence = self.process_image(slide_image_data)

                            # 如果OCR结果有效，使用OCR结果
                            if ocr_text and ocr_text != "图像中未检测到文本" and confidence > 0.3:
                                # 如果有提取的文本，可以结合两者
                                if extracted_text:
                                    final_text = f"{extracted_text}\n\n[OCR识别内容]:\n{ocr_text}"
                                else:
                                    final_text = ocr_text

                                result["metadata"]["slides_extraction_method"][slide_num + 1] = {
                                    "method": "hybrid",
                                    "confidence": confidence
                                }
                                slides_using_ocr += 1
                                logger.info(f"第{slide_num + 1}张幻灯片使用混合处理，OCR置信度: {confidence:.2f}")
                            else:
                                # OCR结果不佳，回退到文本提取
                                final_text = extracted_text
                                result["metadata"]["slides_extraction_method"][slide_num + 1] = {
                                    "method": "text_extraction",
                                    "note": "OCR结果不佳，回退到文本提取"
                                }
                                slides_using_extraction += 1
                                logger.info(f"第{slide_num + 1}张幻灯片OCR结果不佳，使用文本提取")
                        except Exception as e:
                            # OCR处理出错，使用文本提取
                            final_text = extracted_text
                            result["metadata"]["slides_extraction_method"][slide_num + 1] = {
                                "method": "text_extraction",
                                "note": f"OCR处理出错: {str(e)}"
                            }
                            slides_using_extraction += 1
                            logger.error(f"处理第{slide_num + 1}张幻灯片OCR时出错: {str(e)}")
                    else:
                        # 不需要OCR或无法进行OCR，使用文本提取
                        final_text = extracted_text

                        # 记录处理方法和原因
                        note = ""
                        if needs_ocr and not images_available:
                            note = "需要OCR但图像渲染失败"
                        elif needs_ocr and slide_num >= len(slide_images):
                            note = "需要OCR但图像索引超出范围"

                        result["metadata"]["slides_extraction_method"][slide_num + 1] = {
                            "method": "text_extraction",
                            "note": note
                        }
                        slides_using_extraction += 1

                        if note:
                            logger.info(f"第{slide_num + 1}张幻灯片使用文本提取，原因: {note}")
                        else:
                            logger.info(f"第{slide_num + 1}张幻灯片使用文本提取")

                    # 添加到结果中
                    slides_text.append(f"幻灯片 {slide_num + 1}:\n{final_text}")

                # 合并所有幻灯片的文本
                result["text"] = "\n\n".join(slides_text)

                # 添加处理摘要
                result["metadata"]["processing_summary"] = {
                    "slides_using_ocr": slides_using_ocr,
                    "slides_using_extraction": slides_using_extraction,
                    "slides_with_images": slides_with_images,
                    "images_rendering_successful": images_available,
                    "total_images_rendered": len(slide_images) if images_available else 0,
                    "ocr_percentage": slides_using_ocr / total_slides * 100 if total_slides > 0 else 0
                }

                return result

        except Exception as e:
            logger.error(f"PPT处理错误: {str(e)}")
            return {"text": f"处理PPT文档失败: {str(e)}", "metadata": {"error": str(e)}}

    def get_file_content(self, file_data: bytes, file_type: str) -> Dict[str, Any]:
        """
        根据文件类型处理文件并返回文本内容和元数据

        Args:
            file_data: 二进制文件数据
            file_type: MIME类型或文件扩展名

        Returns:
            Dict: 包含提取的文本和处理元数据的字典
        """
        file_type = file_type.lower()

        # 根据文件类型选择处理方法
        if file_type in ['image/jpeg', 'image/png', 'image/jpg', 'image/gif', 'image/bmp', '.jpg', '.jpeg', '.png',
                         '.gif', '.bmp']:
            text, confidence = self.process_image(file_data)
            return {
                "text": text,
                "metadata": {
                    "file_type": "image",
                    "processing_method": "ocr",
                    "confidence": confidence
                }
            }

        elif file_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                           'application/msword', '.docx', '.doc']:
            text = self.process_word_document(file_data)
            return {
                "text": text,
                "metadata": {
                    "file_type": "word",
                    "processing_method": "text_extraction"
                }
            }

        elif file_type in ['application/pdf', '.pdf']:
            # 使用混合PDF处理方法
            return self.process_pdf_hybrid(file_data)

        elif file_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                           'application/vnd.ms-powerpoint', '.pptx', '.ppt']:
            # 判断是否是PPTX格式
            is_pptx = file_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                                    '.pptx']
            # 使用混合PPT处理方法
            return self.process_ppt_hybrid(file_data, is_pptx)

        elif file_type in ['text/html', '.html', '.htm']:
            text = self.sanitize_html(file_data.decode('utf-8', errors='replace'))
            return {
                "text": text,
                "metadata": {
                    "file_type": "html",
                    "processing_method": "html_parsing"
                }
            }

        elif file_type in ['text/plain', '.txt']:
            text = file_data.decode('utf-8', errors='replace')
            return {
                "text": text,
                "metadata": {
                    "file_type": "text",
                    "processing_method": "direct_read"
                }
            }

        else:
            return {
                "text": f"不支持的文件类型: {file_type}",
                "metadata": {
                    "file_type": "unknown",
                    "error": "unsupported_file_type"
                }
            }


# 创建文档处理器实例
doc_processor = DocumentProcessor()


def download_file_from_url(url: str) -> Tuple[bytes, str]:
    """
    从URL下载文件

    Args:
        url: 文件的URL地址

    Returns:
        Tuple[bytes, str]: 文件内容和文件类型
    """
    try:
        response = requests.get(url, timeout=30)
        response.raise_for_status()  # 确保请求成功

        # 获取内容类型
        content_type = response.headers.get('Content-Type', '')

        # 如果无法从header获取类型，尝试从URL推断
        if not content_type or content_type == 'application/octet-stream':
            file_ext = os.path.splitext(url)[1].lower()
            if file_ext:
                # 映射扩展名到MIME类型
                ext_to_mime = {
                    '.pdf': 'application/pdf',
                    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    '.doc': 'application/msword',
                    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    '.ppt': 'application/vnd.ms-powerpoint',
                    '.jpg': 'image/jpeg',
                    '.jpeg': 'image/jpeg',
                    '.png': 'image/png',
                    '.html': 'text/html',
                    '.txt': 'text/plain'
                }
                content_type = ext_to_mime.get(file_ext, 'application/octet-stream')

        return response.content, content_type
    except Exception as e:
        logger.error(f"下载文件失败: {url}, 错误: {str(e)}")
        raise


@app.route('/health', methods=['GET'])
def health_check():
    """健康检查接口"""
    ocr_status = "available" if doc_processor.ocr is not None else "unavailable"
    return jsonify({
        "status": "healthy",
        "ocr_status": ocr_status,
        "version": "1.0.0"
    })


@app.route('/process', methods=['POST'])
def process_file():
    """处理上传的文件"""
    if 'file' not in request.files:
        return jsonify({"error": "请提供文件"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "未选择文件"}), 400

    try:
        # 读取文件数据
        file_data = file.read()
        file_type = file.content_type or os.path.splitext(file.filename)[1]

        # 处理文件
        result = doc_processor.get_file_content(file_data, file_type)

        # 添加文件信息
        result["metadata"]["filename"] = file.filename
        result["metadata"]["original_file_type"] = file_type

        return jsonify(result)
    except Exception as e:
        logger.error(f"处理文件时出错: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/process/pdf', methods=['POST'])
def process_pdf():
    """处理来自URL的PDF文件"""
    if not request.is_json:
        return jsonify({"error": "请提供JSON格式的请求体"}), 400

    data = request.json
    if 'url' not in data:
        return jsonify({"error": "请提供PDF文件URL"}), 400

    try:
        # 下载文件
        file_url = data['url']
        file_data, file_type = download_file_from_url(file_url)

        # 检查文件类型
        if file_type != 'application/pdf' and not file_url.lower().endswith('.pdf'):
            return jsonify({"error": "提供的URL不是PDF文件"}), 400

        # 处理PDF文件
        result = doc_processor.process_pdf_hybrid(file_data)

        # 添加文件信息
        result["metadata"]["file_url"] = file_url
        result["metadata"]["original_file_type"] = file_type

        return jsonify(result)
    except Exception as e:
        logger.error(f"处理URL PDF文件时出错: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/process/ppt', methods=['POST'])
def process_ppt():
    """处理来自URL的PPT文件"""
    if not request.is_json:
        return jsonify({"error": "请提供JSON格式的请求体"}), 400

    data = request.json
    if 'url' not in data:
        return jsonify({"error": "请提供PPT文件URL"}), 400

    try:
        # 下载文件
        file_url = data['url']
        file_data, file_type = download_file_from_url(file_url)

        # 检查文件类型
        is_ppt = file_url.lower().endswith('.ppt') or file_type == 'application/vnd.ms-powerpoint'
        is_pptx = file_url.lower().endswith(
            '.pptx') or file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation'

        if not (is_ppt or is_pptx):
            return jsonify({"error": "提供的URL不是PPT文件"}), 400

        # 处理PPT文件
        result = doc_processor.process_ppt_hybrid(file_data, is_pptx)

        # 添加文件信息
        result["metadata"]["file_url"] = file_url
        result["metadata"]["original_file_type"] = file_type

        return jsonify(result)
    except Exception as e:
        logger.error(f"处理URL PPT文件时出错: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/process/image', methods=['POST'])
def process_image():
    """处理来自URL的图像文件"""
    if not request.is_json:
        return jsonify({"error": "请提供JSON格式的请求体"}), 400

    data = request.json
    if 'url' not in data:
        return jsonify({"error": "请提供图像文件URL"}), 400

    try:
        # 下载文件
        file_url = data['url']
        file_data, file_type = download_file_from_url(file_url)

        # 检查文件类型
        image_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.bmp']
        image_mimetypes = ['image/jpeg', 'image/png', 'image/jpg', 'image/gif', 'image/bmp']

        is_image = any(file_url.lower().endswith(ext) for ext in image_extensions) or file_type in image_mimetypes

        if not is_image:
            return jsonify({"error": "提供的URL不是图像文件"}), 400

        # 处理图像文件
        text, confidence = doc_processor.process_image(file_data)

        result = {
            "text": text,
            "metadata": {
                "file_url": file_url,
                "original_file_type": file_type,
                "file_type": "image",
                "processing_method": "ocr",
                "confidence": confidence
            }
        }

        return jsonify(result)
    except Exception as e:
        logger.error(f"处理URL图像文件时出错: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/process/word', methods=['POST'])
def process_word():
    """处理来自URL的Word文档"""
    if not request.is_json:
        return jsonify({"error": "请提供JSON格式的请求体"}), 400

    data = request.json
    if 'url' not in data:
        return jsonify({"error": "请提供Word文档URL"}), 400

    try:
        # 下载文件
        file_url = data['url']
        file_data, file_type = download_file_from_url(file_url)

        # 检查文件类型
        word_extensions = ['.doc', '.docx']
        word_mimetypes = [
            'application/msword',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        ]

        is_word = any(file_url.lower().endswith(ext) for ext in word_extensions) or file_type in word_mimetypes

        if not is_word:
            return jsonify({"error": "提供的URL不是Word文档"}), 400

        # 处理Word文档
        text = doc_processor.process_word_document(file_data)

        result = {
            "text": text,
            "metadata": {
                "file_url": file_url,
                "original_file_type": file_type,
                "file_type": "word",
                "processing_method": "text_extraction"
            }
        }

        return jsonify(result)
    except Exception as e:
        logger.error(f"处理URL Word文档时出错: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/process/html', methods=['POST'])
def process_html():
    """处理HTML内容"""
    if not request.is_json:
        return jsonify({"error": "请提供JSON格式的请求体"}), 400

    data = request.json
    if 'html' not in data:
        return jsonify({"error": "请提供HTML内容"}), 400

    try:
        html_content = data['html']
        sanitized_text = doc_processor.sanitize_html(html_content)

        result = {
            "text": sanitized_text,
            "metadata": {
                "file_type": "html",
                "processing_method": "html_parsing",
                "original_length": len(html_content),
                "processed_length": len(sanitized_text)
            }
        }

        return jsonify(result)
    except Exception as e:
        logger.error(f"处理HTML内容时出错: {str(e)}")
        return jsonify({"error": str(e)}), 500


if __name__ == "__main__":
    # 如果存在PORT环境变量，使用它，否则使用默认的5002
    port = int(os.environ.get("PORT", 5002))
    # 在开发模式下启用调试，生产环境应使用 gunicorn 或 uwsgi